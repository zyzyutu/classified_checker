# -*- coding: utf-8 -*-
"""
文件检查模块 - 支持多种文档格式和压缩包的涉密信息检查
包含：TXT、DOC、DOCX、XLS、XLSX、PPT、PPTX、PDF、ZIP、RAR、7Z
功能：Magic Number校验、加密文件识别、压缩包递归解压、多线程并行检查
"""

import os
import tempfile
import shutil
import threading
from collections import Counter
from concurrent.futures import ThreadPoolExecutor, as_completed

from utils import build_combined_pattern

# ========== Magic Number 定义（文件头校验真实类型） ==========
MAGIC_NUMBERS = {
    '.doc':  [b'\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1'],          # OLE2格式
    '.docx': [b'PK\x03\x04'],                                    # ZIP格式
    '.xls':  [b'\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1'],          # OLE2格式
    '.xlsx': [b'PK\x03\x04'],                                    # ZIP格式
    '.ppt':  [b'\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1'],          # OLE2格式
    '.pptx': [b'PK\x03\x04'],                                    # ZIP格式
    '.pdf':  [b'%PDF'],                                           # PDF格式
    '.txt':  None,  # TXT无固定Magic Number，不校验
    '.zip':  [b'PK\x03\x04'],
    '.rar':  [b'Rar!\x1a\x07'],
    '.7z':   [b'7z\xbc\xaf\x27\x1c'],
}

ARCHIVE_EXTS = {'.zip', '.rar', '.7z'}
DOC_EXTS = {'.doc', '.xls', '.ppt'}  # 需要 COM 接口的文件类型
MAX_ARCHIVE_DEPTH = 3  # 压缩包最大嵌套深度

# COM 线程锁（WPS/Word COM 不是线程安全的）
_com_lock = threading.Lock()


# ==================== 压缩包处理 ====================

def _is_archive(fpath):
    """判断文件是否为支持的压缩包"""
    ext = os.path.splitext(fpath)[1].lower()
    if ext not in ARCHIVE_EXTS:
        return False
    try:
        with open(fpath, 'rb') as f:
            header = f.read(8)
    except Exception:
        return False
    if ext == '.zip':
        return header.startswith(b'PK\x03\x04')
    if ext == '.rar':
        return header.startswith(b'Rar!\x1a\x07')
    if ext == '.7z':
        return header.startswith(b'7z\xbc\xaf\x27\x1c')
    return False


def _extract_archive(fpath, extract_dir):
    """
    解压压缩包到指定目录。
    返回: (extracted_files, is_encrypted)
    """
    ext = os.path.splitext(fpath)[1].lower()
    extracted = []

    try:
        if ext == '.zip':
            import zipfile
            try:
                zf = zipfile.ZipFile(fpath, 'r')
            except zipfile.BadZipFile:
                return [], False
            except RuntimeError:
                # 密码保护的 ZIP
                return [], True
            with zf:
                # 检查是否有加密文件
                has_encrypted = any(info.flag_bits & 0x1 for info in zf.infolist())
                if has_encrypted:
                    return [], True
                for info in zf.infolist():
                    if info.is_dir():
                        continue
                    basename = os.path.basename(info.filename)
                    if basename.startswith('.') or basename.startswith('__'):
                        continue
                    target = os.path.join(extract_dir, info.filename)
                    os.makedirs(os.path.dirname(target), exist_ok=True)
                    with zf.open(info) as src, open(target, 'wb') as dst:
                        dst.write(src.read())
                    extracted.append(target)

        elif ext == '.rar':
            try:
                import rarfile
                with rarfile.RarFile(fpath) as rf:
                    # 检查是否需要密码
                    if rf.needs_password():
                        return [], True
                    for info in rf.infolist():
                        if info.is_dir():
                            continue
                        basename = os.path.basename(info.filename)
                        if basename.startswith('.'):
                            continue
                        target = os.path.join(extract_dir, info.filename)
                        os.makedirs(os.path.dirname(target), exist_ok=True)
                        with rf.open(info) as src, open(target, 'wb') as dst:
                            dst.write(src.read())
                        extracted.append(target)
            except ImportError:
                pass
            except rarfile.BadRarFile:
                return [], False

        elif ext == '.7z':
            try:
                import py7zr
                try:
                    sz = py7zr.SevenZipFile(fpath, mode='r')
                except py7zr.Bad7zFile:
                    return [], False
                with sz:
                    if sz.needs_password():
                        return [], True
                    sz.extractall(path=extract_dir)
                for root, dirs, fnames in os.walk(extract_dir):
                    for fname in fnames:
                        extracted.append(os.path.join(root, fname))
            except ImportError:
                pass
            except Exception:
                return [], False

    except Exception:
        return [], False

    return extracted, False


def _check_archive(fpath, pattern, log_callback, depth, visited_archives):
    """
    解压并递归检查压缩包内的文件。

    返回: (details, is_encrypted)
    """
    if depth > MAX_ARCHIVE_DEPTH:
        if log_callback:
            log_callback(f"  [文件] 压缩包嵌套过深，跳过: {fpath}")
        return [], False

    if fpath in visited_archives:
        return [], False
    visited_archives.add(fpath)

    tmp_dir = tempfile.mkdtemp(prefix="checker_")
    details = []

    try:
        extracted, is_encrypted = _extract_archive(fpath, tmp_dir)
        if is_encrypted:
            return [], True
        if not extracted:
            if log_callback:
                log_callback(f"  [文件] 压缩包解压失败或为空: {fpath}")
            return [], False

        if log_callback:
            log_callback(f"  [文件] 压缩包解压出 {len(extracted)} 个文件: {fpath}")

        for epath in extracted:
            ext = os.path.splitext(epath)[1].lower()

            # 嵌套压缩包递归处理
            if ext in ARCHIVE_EXTS and _is_archive(epath):
                sub_details, _ = _check_archive(
                    epath, pattern, log_callback, depth + 1, visited_archives)
                details.extend(sub_details)
                continue

            # 普通文件：检查内容
            if ext not in MAGIC_NUMBERS:
                continue

            try:
                file_text = _extract_text(epath, ext)
                if file_text is None:
                    continue
                if file_text.startswith("[") and "加密" in file_text:
                    continue

                lines = file_text.split("\n")
                for i, line in enumerate(lines, 1):
                    line_stripped = line.strip()
                    if not line_stripped:
                        continue
                    for m in pattern.finditer(line_stripped):
                        details.append({
                            "file": fpath,
                            "line_no": i,
                            "content": line_stripped[:120],
                            "keyword": m.group(),
                            "file_type": f"[压缩包内{ext}]"
                        })
            except Exception:
                pass

    finally:
        try:
            shutil.rmtree(tmp_dir, ignore_errors=True)
        except Exception:
            pass

    return details, False


# ==================== 主检查函数 ====================

def check_files(directory, keywords, log_callback=None, max_workers=6):
    """
    递归扫描目录下所有支持的文件和压缩包，多线程并行检查涉密信息。

    参数:
        directory:    扫描目录路径
        keywords:     关键词列表
        log_callback: 日志回调函数（可选）
        max_workers:  并行线程数（默认6）

    返回:
        dict: {
            "total_files": 扫描文件总数,
            "supported_files": 支持格式文件数,
            "matched_files": 涉密文件数,
            "archives_scanned": 扫描压缩包数,
            "type_counts": {ext: count},
            "encrypted_files": [路径],
            "hidden_files": [路径],
            "details": [{file, line_no, content, keyword, file_type}, ...]
        }
    """
    pattern = build_combined_pattern(keywords)
    empty_result = {
        "total_files": 0, "supported_files": 0, "matched_files": 0,
        "type_counts": {}, "encrypted_files": [], "damaged_files": [],
        "hidden_files": [], "archives_scanned": 0,
        "encrypted_archives": [], "details": []
    }
    if not pattern:
        return empty_result
    if not os.path.isdir(directory):
        if log_callback:
            log_callback(f"  [文件] 目录不存在: {directory}")
        return empty_result

    # ========== 阶段一：收集所有文件 ==========
    all_files = []
    hidden_files = []

    for root, dirs, files in os.walk(directory):
        for fname in files:
            fpath = os.path.join(root, fname)
            ext = os.path.splitext(fname)[1].lower()

            # 检测隐藏文件（Windows属性 + Linux点前缀）
            is_hidden = False
            try:
                import ctypes
                attrs = ctypes.windll.kernel32.GetFileAttributesW(fpath)
                if attrs != -1 and (attrs & 0x02):
                    is_hidden = True
            except Exception:
                pass
            if fname.startswith("."):
                is_hidden = True
            if is_hidden:
                hidden_files.append(fpath)
                if log_callback:
                    log_callback(f"  [文件] 发现隐藏文件: {fpath}")

            # 收集支持的普通文件
            if ext in MAGIC_NUMBERS and ext not in ARCHIVE_EXTS:
                all_files.append((fpath, ext, False))
            # 收集压缩包
            elif ext in ARCHIVE_EXTS:
                all_files.append((fpath, ext, True))

    total_files = len([f for f in all_files if not f[2]])
    archive_files = [f for f in all_files if f[2]]

    if log_callback:
        log_callback(f"  [文件] 收集完成: {total_files} 个普通文件, "
                     f"{len(archive_files)} 个压缩包")

    # ========== 阶段二：先处理压缩包（顺序，因为涉及临时目录） ==========
    matched_files = set()
    details = []
    type_counts = Counter()
    encrypted_files = []
    damaged_files = []
    archives_scanned = 0
    encrypted_archives = []
    visited_archives = set()

    for fpath, ext, _ in archive_files:
        archives_scanned += 1
        if log_callback:
            log_callback(f"  [文件] 正在解压检查压缩包: {fpath}")

        archive_details, archive_encrypted = _check_archive(
            fpath, pattern, log_callback, 0, visited_archives)
        if archive_encrypted:
            encrypted_archives.append(fpath)
            encrypted_files.append(fpath)
            if log_callback:
                log_callback(f"  [文件] 压缩包需要密码: {fpath}")
        for d in archive_details:
            matched_files.add(d["file"])
        details.extend(archive_details)

    # ========== 阶段三：多线程并行检查普通文件 ==========
    def process_one_file(fpath, ext):
        """处理单个文件，返回 (result_type, data)"""
        result_type = "normal"
        data = {
            "file": fpath,
            "ext": ext,
            "matched": False,
            "encrypted": False,
            "damaged": False,
            "status_msg": "",
            "details": []
        }

        if log_callback:
            log_callback(f"  [文件] 正在检查: {fpath}")

        # Magic Number 校验
        real_ext = _check_magic_number(fpath, ext)
        if real_ext != ext:
            if log_callback:
                log_callback(f"  [文件] 类型不匹配: {os.path.basename(fpath)} "
                             f"(扩展名{ext}, 实际{real_ext})")

        # 提取文本，区分加密/损坏/格式不支持
        file_text = None
        try:
            if ext in DOC_EXTS:
                with _com_lock:
                    file_text = _extract_text(fpath, ext)
            else:
                file_text = _extract_text(fpath, ext)
        except Exception as e:
            data["damaged"] = True
            data["status_msg"] = f"读取异常: {e}"
            return result_type, data

        # 判断提取结果
        if file_text is None:
            # 进一步区分：是加密还是损坏还是库未安装
            reason = _diagnose_extract_failure(fpath, ext)
            if reason == "encrypted":
                data["encrypted"] = True
                data["status_msg"] = "文件已加密"
            elif reason == "damaged":
                data["damaged"] = True
                data["status_msg"] = "文件损坏或格式异常"
            elif reason == "missing_lib":
                data["status_msg"] = "缺少解析库"
            else:
                data["damaged"] = True
                data["status_msg"] = "无法读取"
            return result_type, data

        if file_text.startswith("[") and "加密" in file_text:
            data["encrypted"] = True
            data["status_msg"] = file_text.strip("[]")
            return result_type, data

        if file_text.startswith("[") and "损坏" in file_text:
            data["damaged"] = True
            data["status_msg"] = file_text.strip("[]")
            return result_type, data

        # 检查关键词
        lines = file_text.split("\n")
        file_details = []
        for i, line in enumerate(lines, 1):
            line_stripped = line.strip()
            if not line_stripped:
                continue
            for m in pattern.finditer(line_stripped):
                file_details.append({
                    "file": fpath,
                    "line_no": i,
                    "content": line_stripped[:120],
                    "keyword": m.group(),
                    "file_type": ext
                })

        if file_details:
            data["matched"] = True
            data["details"] = file_details

        return result_type, data

    if log_callback:
        normal_count = len([f for f in all_files if not f[2]])
        log_callback(f"  [文件] 开始并行检查 {normal_count} 个文件 (线程{max_workers})")

    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        futures = {}
        for fpath, ext, is_archive in all_files:
            if not is_archive:
                futures[executor.submit(process_one_file, fpath, ext)] = fpath

        for future in as_completed(futures):
            try:
                _, data = future.result()
                if data["encrypted"]:
                    encrypted_files.append(data["file"])
                    if log_callback:
                        log_callback(f"  [文件] 已加密: {data['file']} "
                                     f"({data.get('status_msg', '')})")
                elif data["damaged"]:
                    damaged_files.append(data["file"])
                    if log_callback:
                        log_callback(f"  [文件] 文件损坏: {data['file']} "
                                     f"({data.get('status_msg', '')})")
                elif data["matched"]:
                    matched_files.add(data["file"])
                    details.extend(data["details"])
                type_counts[data["ext"]] += 1
            except Exception as e:
                if log_callback:
                    log_callback(f"  [文件] 处理异常: {futures[future]} - {e}")

    return {
        "total_files": total_files,
        "supported_files": total_files,
        "matched_files": len(matched_files),
        "type_counts": dict(type_counts),
        "encrypted_files": encrypted_files,
        "damaged_files": damaged_files,
        "hidden_files": hidden_files,
        "archives_scanned": archives_scanned,
        "encrypted_archives": encrypted_archives,
        "details": details
    }


# ==================== 辅助函数 ====================

def _diagnose_extract_failure(fpath, ext):
    """
    诊断文本提取失败的原因。
    返回: "encrypted" | "damaged" | "missing_lib" | "unknown"
    """
    try:
        if ext == '.docx':
            import zipfile
            if not zipfile.is_zipfile(fpath):
                return "damaged"
            with zipfile.ZipFile(fpath, 'r') as zf:
                # 检查是否加密（ZIP加密标志位）
                for info in zf.infolist():
                    if info.flag_bits & 0x1:
                        return "encrypted"
            # ZIP有效但python-docx打开失败 → 可能损坏
            return "damaged"

        elif ext == '.xlsx':
            import zipfile
            if not zipfile.is_zipfile(fpath):
                return "damaged"
            with zipfile.ZipFile(fpath, 'r') as zf:
                for info in zf.infolist():
                    if info.flag_bits & 0x1:
                        return "encrypted"
            return "damaged"

        elif ext == '.pptx':
            import zipfile
            if not zipfile.is_zipfile(fpath):
                return "damaged"
            with zipfile.ZipFile(fpath, 'r') as zf:
                for info in zf.infolist():
                    if info.flag_bits & 0x1:
                        return "encrypted"
            return "damaged"

        elif ext == '.xls':
            try:
                import xlrd
                xlrd.open_workbook(fpath)
            except xlrd.XLRDError as e:
                if "password" in str(e).lower() or "encrypt" in str(e).lower():
                    return "encrypted"
                return "damaged"
            except ImportError:
                return "missing_lib"
            return "damaged"

        elif ext == '.pdf':
            try:
                import PyPDF2
                with open(fpath, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    if reader.is_encrypted:
                        return "encrypted"
            except ImportError:
                return "missing_lib"
            return "damaged"

    except ImportError:
        return "missing_lib"
    except Exception:
        return "unknown"

    return "unknown"


def _check_magic_number(fpath, expected_ext):
    """通过文件头Magic Number校验文件真实类型"""
    try:
        with open(fpath, 'rb') as f:
            header = f.read(16)
    except Exception:
        return expected_ext

    if header.startswith(b'PK\x03\x04'):
        try:
            import zipfile
            if zipfile.is_zipfile(fpath):
                with zipfile.ZipFile(fpath, 'r') as zf:
                    names = zf.namelist()
                    if any('word/' in n for n in names):
                        return '.docx'
                    elif any('xl/' in n for n in names):
                        return '.xlsx'
                    elif any('ppt/' in n for n in names):
                        return '.pptx'
        except Exception:
            pass
        return '.docx'

    if header.startswith(b'\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1'):
        return expected_ext

    if header.startswith(b'%PDF'):
        return '.pdf'

    return expected_ext


def _extract_text(fpath, ext):
    """根据文件类型提取文本内容"""
    extractors = {
        '.txt': _extract_txt,
        '.doc': _extract_doc,
        '.docx': _extract_docx,
        '.xls': _extract_xls,
        '.xlsx': _extract_xlsx,
        '.ppt': _extract_ppt,
        '.pptx': _extract_pptx,
        '.pdf': _extract_pdf,
    }
    extractor = extractors.get(ext)
    if extractor:
        return extractor(fpath)
    return None


def _extract_txt(fpath):
    for enc in ('utf-8', 'gbk', 'gb2312', 'latin-1'):
        try:
            with open(fpath, 'r', encoding=enc) as f:
                return f.read()
        except (UnicodeDecodeError, UnicodeError):
            continue
    return None


def _extract_doc(fpath):
    """提取旧版DOC文件文本（OLE2格式），使用 WPS/Word COM 接口"""
    try:
        import pythoncom
        import win32com.client

        pythoncom.CoInitialize()
        app = None
        doc = None
        try:
            for prog_id in ("KWps.Application", "WPS.Application",
                            "Word.Application", "et.Application"):
                try:
                    app = win32com.client.Dispatch(prog_id)
                    app.Visible = False
                    app.DisplayAlerts = False
                    break
                except Exception:
                    app = None
                    continue

            if app is None:
                return None

            abs_path = os.path.abspath(fpath)
            doc = app.Documents.Open(abs_path, ReadOnly=True,
                                     PasswordDocument="",
                                     PasswordTemplate="",
                                     NoEncodingDialog=True)
            full_text = doc.Content.Text
            return full_text if full_text else None
        finally:
            if doc is not None:
                try:
                    doc.Close(SaveChanges=False)
                except Exception:
                    pass
            if app is not None:
                try:
                    app.Quit()
                except Exception:
                    pass
            pythoncom.CoUninitialize()
    except ImportError:
        pass
    except Exception:
        pass

    try:
        import olefile
        if olefile.isOleFile(fpath):
            ole = olefile.OleFileIO(fpath)
            is_encrypted = ole.exists('EncryptionInfo')
            ole.close()
            if is_encrypted:
                return "[DOC文件-已加密，无法读取]"
    except Exception:
        pass

    try:
        with open(fpath, 'rb') as f:
            raw = f.read()
        import re
        text_parts = []
        i = 0
        while i < len(raw):
            b = raw[i]
            if 0xB0 <= b <= 0xF7 and i + 1 < len(raw):
                b2 = raw[i + 1]
                if 0xA1 <= b2 <= 0xFE:
                    start = i
                    while i < len(raw) - 1:
                        b = raw[i]
                        b2 = raw[i + 1]
                        if 0xB0 <= b <= 0xF7 and 0xA1 <= b2 <= 0xFE:
                            i += 2
                        else:
                            break
                    try:
                        segment = raw[start:i].decode('gbk', errors='ignore')
                        if len(segment) >= 2:
                            text_parts.append(segment)
                    except Exception:
                        pass
                    continue
            if 0x20 <= b <= 0x7E:
                start = i
                while i < len(raw) and 0x20 <= raw[i] <= 0x7E:
                    i += 1
                segment = raw[start:i].decode('ascii', errors='ignore')
                if len(segment) >= 4:
                    text_parts.append(segment)
                continue
            i += 1
        result = "\n".join(text_parts)
        return result if result else None
    except Exception:
        return None


def _extract_docx(fpath):
    try:
        from docx import Document
        doc = Document(fpath)
        paragraphs = [p.text for p in doc.paragraphs]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    paragraphs.append(cell.text)
        return "\n".join(paragraphs)
    except Exception:
        return None


def _extract_xls(fpath):
    try:
        import xlrd
        wb = xlrd.open_workbook(fpath)
        texts = []
        for sheet in wb.sheets():
            for row_idx in range(sheet.nrows):
                row_values = [str(sheet.cell_value(row_idx, col))
                              for col in range(sheet.ncols)]
                texts.append("\t".join(row_values))
        return "\n".join(texts)
    except Exception:
        return None


def _extract_xlsx(fpath):
    try:
        from openpyxl import load_workbook
        wb = load_workbook(fpath, data_only=True)
        texts = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join(str(c) if c is not None else ""
                                     for c in row)
                texts.append(row_text)
        return "\n".join(texts)
    except Exception:
        return None


def _extract_ppt(fpath):
    try:
        import pythoncom
        import win32com.client

        pythoncom.CoInitialize()
        app = None
        pres = None
        try:
            for prog_id in ("KWpp.Application", "WPP.Application",
                            "PowerPoint.Application"):
                try:
                    app = win32com.client.Dispatch(prog_id)
                    app.Visible = False
                    break
                except Exception:
                    app = None
                    continue

            if app is None:
                return None

            abs_path = os.path.abspath(fpath)
            pres = app.Presentations.Open(abs_path, ReadOnly=True,
                                          WithWindow=False)
            texts = []
            for slide in pres.Slides:
                for shape in slide.Shapes:
                    if shape.HasTextFrame:
                        texts.append(shape.TextFrame.TextRange.Text)
                    if shape.HasTable:
                        for row in shape.Table.Rows:
                            for cell in row.Cells:
                                texts.append(cell.Shape.TextFrame
                                             .TextRange.Text)
            return "\n".join(texts) if texts else None
        finally:
            if pres is not None:
                try:
                    pres.Close()
                except Exception:
                    pass
            if app is not None:
                try:
                    app.Quit()
                except Exception:
                    pass
            pythoncom.CoUninitialize()
    except ImportError:
        pass
    except Exception:
        pass
    return None


def _extract_pptx(fpath):
    try:
        from pptx import Presentation
        prs = Presentation(fpath)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        texts.append(paragraph.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            texts.append(cell.text)
        return "\n".join(texts)
    except Exception:
        return None


def _extract_pdf(fpath):
    try:
        import PyPDF2
        texts = []
        with open(fpath, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            if reader.is_encrypted:
                return None
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    texts.append(text)
        return "\n".join(texts)
    except Exception:
        pass

    try:
        import pdfplumber
        texts = []
        with pdfplumber.open(fpath) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    texts.append(text)
        return "\n".join(texts)
    except Exception:
        return None
