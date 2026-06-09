# -*- coding: utf-8 -*-
"""
文本提取模块 - 从各种文件格式中提取纯文本
支持：TXT、DOC、DOCX、XLS、XLSX、PPT、PPTX、PDF
包含快速提取器（calamine/pypdfium2）和标准提取器两级机制
"""

import os
import threading
from contextlib import contextmanager

# COM 线程锁（WPS/Word COM 不是线程安全的，RLock允许同一线程重入）
_com_lock = threading.RLock()
_thread_local = threading.local()

# COM 应用实例缓存（避免每次文件都重启 Word/WPS）
_cached_app = None
_cached_prog_id = None


# ==================== COM 公共工具 ====================

@contextmanager
def _com_app(prog_ids):
    """
    COM 应用程序上下文管理器（带缓存）。
    首次调用启动应用并缓存，后续调用直接复用，避免每次重启的 2-5 秒开销。
    注意：需要配合 RLock 使用（同一线程可能重入）。
    """
    global _cached_app, _cached_prog_id
    import pythoncom
    import win32com.client

    if not hasattr(_thread_local, 'com_initialized'):
        pythoncom.CoInitialize()
        _thread_local.com_initialized = True

    # 命中缓存：直接复用
    if _cached_app is not None:
        try:
            _cached_app.Name  # 探活：访问一个属性确认应用还活着
            yield _cached_app
            return
        except Exception:
            _cached_app = None
            _cached_prog_id = None

    # 未命中：启动新实例
    app = None
    for prog_id in prog_ids:
        try:
            app = win32com.client.dynamic.Dispatch(prog_id)
            if hasattr(app, 'DisplayAlerts'):
                app.DisplayAlerts = False
            break
        except Exception:
            app = None
            continue

    if app is None:
        raise RuntimeError("无法启动COM应用程序")

    _cached_app = app
    _cached_prog_id = prog_ids[0]
    try:
        yield app
    except Exception:
        _cached_app = None
        _cached_prog_id = None
        raise


def _com_app_quit():
    """释放缓存的 COM 应用（检查全部完成后调用）"""
    global _cached_app, _cached_prog_id
    if _cached_app is not None:
        try:
            _cached_app.Quit()
        except Exception:
            pass
        _cached_app = None
        _cached_prog_id = None


def _try_com_extract(fpath, prog_ids, extract_func):
    """
    安全执行 COM 提取：加锁 → 调用 _com_app → 提取 → 释放。
    extract_func(app, abs_path) -> str or None
    """
    with _com_lock:
        try:
            with _com_app(prog_ids) as app:
                return extract_func(app, os.path.abspath(fpath))
        except RuntimeError:
            return None  # COM 启动失败
        except ImportError:
            return None  # 缺少 pywin32
        except Exception:
            return None  # 其他错误


# ==================== 快速提取器 ====================

def _extract_xlsx_fast(fpath):
    """calamine（Rust实现，比openpyxl快10-50倍）"""
    try:
        from python_calamine import load_workbook
        wb = load_workbook(fpath)
        texts = []
        for sheet_name in wb.sheet_names:
            for row in wb.get_sheet_by_name(sheet_name).rows:
                texts.append("\t".join(str(c) if c is not None else "" for c in row))
        return "\n".join(texts) if texts else None
    except (ImportError, Exception):
        return None


def _extract_pdf_fast(fpath):
    """pypdfium2（Google维护，比PyPDF2快3-5倍）"""
    try:
        import pypdfium2 as pdfium
        pdf = pdfium.PdfDocument(fpath)
        texts = []
        for i in range(len(pdf)):
            text = pdf[i].get_text()
            if text:
                texts.append(text)
        return "\n".join(texts) if texts else None
    except (ImportError, Exception):
        return None


# ==================== 标准提取器 ====================

def _extract_txt(fpath):
    """TXT：尝试多种编码"""
    for enc in ('utf-8', 'gbk', 'gb2312', 'latin-1'):
        try:
            with open(fpath, 'r', encoding=enc) as f:
                return f.read()
        except (UnicodeDecodeError, UnicodeError):
            continue
    return None


def _extract_doc(fpath):
    """DOC：COM接口 + olefile加密检测 + raw字节兜底"""

    def _com_open(app, abs_path):
        doc = app.Documents.Open(abs_path, ReadOnly=True,
                                 PasswordDocument="", PasswordTemplate="",
                                 NoEncodingDialog=True)
        try:
            return doc.Content.Text or None
        finally:
            doc.Close(SaveChanges=False)

    result = _try_com_extract(fpath, ("KWps.Application", "WPS.Application",
                                      "Word.Application"), _com_open)
    if result:
        return result

    # COM 失败：检测是否加密
    try:
        reason = diagnose_extract_failure(fpath, '.doc')
        if reason == "encrypted":
            return "[DOC文件-已加密，无法读取]"
    except Exception:
        pass

    # raw字节兜底：提取中文连续段+ASCII段
    return _extract_raw_text(fpath)


def _extract_docx(fpath):
    """DOCX：python-docx"""
    try:
        from docx import Document
        doc = Document(fpath)
        parts = [p.text for p in doc.paragraphs]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    parts.append(cell.text)
        return "\n".join(parts)
    except Exception:
        return None


def _extract_xls(fpath):
    """XLS：xlrd"""
    try:
        import xlrd
        wb = xlrd.open_workbook(fpath)
        texts = []
        for sheet in wb.sheets():
            for r in range(sheet.nrows):
                texts.append("\t".join(str(sheet.cell_value(r, c))
                                       for c in range(sheet.ncols)))
        return "\n".join(texts)
    except Exception:
        return None


def _extract_xlsx(fpath):
    """XLSX：openpyxl（标准方法）"""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(fpath, data_only=True)
        texts = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                texts.append("\t".join(str(c) if c is not None else "" for c in row))
        return "\n".join(texts)
    except Exception:
        return None


def _extract_ppt(fpath):
    """PPT：COM接口提取幻灯片文本"""

    def _com_open(app, abs_path):
        pres = app.Presentations.Open(abs_path, WithWindow=False, ReadOnly=True)
        try:
            texts = []
            for slide in pres.Slides:
                for shape in slide.Shapes:
                    try:
                        if shape.HasTextFrame:
                            texts.append(shape.TextFrame.TextRange.Text)
                        if hasattr(shape, 'HasTable') and shape.HasTable:
                            for row in shape.Table.Rows:
                                for cell in row.Cells:
                                    texts.append(cell.Shape.TextFrame.TextRange.Text)
                    except Exception:
                        continue
            return "\n".join(texts) if texts else None
        finally:
            pres.Close()

    result = _try_com_extract(fpath, ("KWps.Application", "WPP.Application",
                                      "PowerPoint.Application"), _com_open)
    return result or "[PPT文件-无法打开，请安装WPS或PowerPoint]"


def _extract_pptx(fpath):
    """PPTX：python-pptx"""
    try:
        from pptx import Presentation
        prs = Presentation(fpath)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
        return "\n".join(texts)
    except Exception:
        return None


def _extract_pdf(fpath):
    """PDF：PyPDF2优先，pdfplumber降级"""
    try:
        import PyPDF2
        with open(fpath, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            if reader.is_encrypted:
                return None
            texts = [page.extract_text() for page in reader.pages]
            return "\n".join(t for t in texts if t)
    except Exception:
        pass

    try:
        import pdfplumber
        with pdfplumber.open(fpath) as pdf:
            texts = [page.extract_text() for page in pdf.pages]
            return "\n".join(t for t in texts if t)
    except Exception:
        return None


# ==================== 辅助函数 ====================

def _extract_raw_text(fpath):
    """从二进制文件中提取中文连续段+ASCII段（兜底方案）"""
    try:
        with open(fpath, 'rb') as f:
            raw = f.read()
        parts = []
        i = 0
        while i < len(raw):
            b = raw[i]
            # GBK双字节中文
            if 0xB0 <= b <= 0xF7 and i + 1 < len(raw):
                b2 = raw[i + 1]
                if 0xA1 <= b2 <= 0xFE:
                    start = i
                    while i < len(raw) - 1:
                        if 0xB0 <= raw[i] <= 0xF7 and 0xA1 <= raw[i+1] <= 0xFE:
                            i += 2
                        else:
                            break
                    seg = raw[start:i].decode('gbk', errors='ignore')
                    if len(seg) >= 2:
                        parts.append(seg)
                    continue
            # ASCII可打印字符
            if 0x20 <= b <= 0x7E:
                start = i
                while i < len(raw) and 0x20 <= raw[i] <= 0x7E:
                    i += 1
                seg = raw[start:i].decode('ascii', errors='ignore')
                if len(seg) >= 4:
                    parts.append(seg)
                continue
            i += 1
        return "\n".join(parts) if parts else None
    except Exception:
        return None


# ==================== 分发器 ====================

# 快速提取器（优先尝试，失败自动降级）
_FAST = {
    '.xlsx': _extract_xlsx_fast,
    '.pdf': _extract_pdf_fast,
}

# 标准提取器（降级方案）
_STANDARD = {
    '.txt': _extract_txt,
    '.doc': _extract_doc,
    '.docx': _extract_docx,
    '.xls': _extract_xls,
    '.xlsx': _extract_xlsx,
    '.ppt': _extract_ppt,
    '.pptx': _extract_pptx,
    '.pdf': _extract_pdf,
}


def extract_text(fpath, ext):
    """
    根据文件类型提取文本内容。
    快速提取器优先，失败自动降级到标准提取器。
    """
    fast = _FAST.get(ext)
    if fast:
        try:
            result = fast(fpath)
            if result is not None:
                return result
        except Exception:
            pass

    extractor = _STANDARD.get(ext)
    return extractor(fpath) if extractor else None


def diagnose_extract_failure(fpath, ext):
    """
    诊断文本提取失败的原因。
    返回: "encrypted" | "damaged" | "missing_lib" | "unknown"
    """
    try:
        # ZIP格式（docx/xlsx/pptx）：多层检测
        if ext in ('.docx', '.xlsx', '.pptx'):
            import zipfile
            if not zipfile.is_zipfile(fpath):
                return "damaged"
            with zipfile.ZipFile(fpath, 'r') as zf:
                names = zf.namelist()
                # 1. 标准ZIP加密位
                if any(info.flag_bits & 0x1 for info in zf.infolist()):
                    return "encrypted"
                # 2. WPS/Office加密流（EncryptedPackage）
                if any('EncryptedPackage' in n for n in names):
                    return "encrypted"
                # 3. 应用层加密：ZIP合法但缺少核心内容流（只有Content_Types.xml）
                core_streams = {
                    '.docx': 'word/document.xml',
                    '.xlsx': 'xl/sharedStrings.xml',
                    '.pptx': 'ppt/slides/slide1.xml',
                }
                expected = core_streams.get(ext)
                if expected and not any(expected in n for n in names):
                    return "encrypted"
            return "damaged"

        # OLE2格式（doc/xls/ppt）：检查加密流
        if ext in ('.doc', '.xls', '.ppt'):
            try:
                import olefile
                if not olefile.isOleFile(fpath):
                    return "damaged"
                ole = olefile.OleFileIO(fpath)
                encrypted = ole.exists('EncryptionInfo')
                ole.close()
                return "encrypted" if encrypted else "missing_lib"
            except ImportError:
                return "missing_lib"
            except Exception:
                return "damaged"

        # PDF：检查加密属性
        if ext == '.pdf':
            try:
                import PyPDF2
                with open(fpath, 'rb') as f:
                    if PyPDF2.PdfReader(f).is_encrypted:
                        return "encrypted"
            except ImportError:
                return "missing_lib"
            except Exception:
                return "damaged"
            return "damaged"

    except ImportError:
        return "missing_lib"
    except Exception:
        pass

    return "unknown"
